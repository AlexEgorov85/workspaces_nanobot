from __future__ import annotations

import csv
import mimetypes
from pathlib import Path

import chardet


def detect_format(path: str | Path) -> str:
    p = Path(path)
    ext = p.suffix.lower().lstrip(".")
    if ext:
        return ext
    guess, _ = mimetypes.guess_type(str(p))
    if guess:
        return guess.split("/")[-1]
    return ""


def _read_text_auto(path: Path) -> str:
    raw = path.read_bytes()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding") or "utf-8"
    confidence = float(detected.get("confidence") or 0.0)
    if confidence < 0.7:
        for fallback in ("utf-8", "cp1251", "latin-1"):
            try:
                return raw.decode(fallback)
            except UnicodeDecodeError:
                continue
        return raw.decode("utf-8", errors="replace")
    return raw.decode(encoding, errors="replace")


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    pres = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(pres.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        chunks.append(txt)
        if chunks:
            parts.append(f"--- слайд {idx} ---\n" + "\n".join(chunks))
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- лист: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells))
    wb.close()
    return "\n\n".join(parts)


def _extract_xls(path: Path) -> str:
    import xlrd

    book = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"--- лист: {sheet.name} ---")
        for row_idx in range(sheet.nrows):
            row = sheet.row_values(row_idx)
            cells = [str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_csv(path: Path) -> str:
    text = _read_text_auto(path)
    reader = csv.reader(text.splitlines())
    rows = [" | ".join(row) for row in reader if row]
    return "\n".join(rows)


def extract_text(path: str | Path) -> str:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"Файл не найден: {p}")
    fmt = detect_format(p)
    if fmt in ("docx",):
        return _extract_docx(p)
    if fmt in ("xlsx",):
        return _extract_xlsx(p)
    if fmt in ("xls",):
        return _extract_xls(p)
    if fmt in ("pdf",):
        return _extract_pdf(p)
    if fmt in ("pptx",):
        return _extract_pptx(p)
    if fmt in ("csv",):
        return _extract_csv(p)
    if fmt in ("txt",):
        return _read_text_auto(p)
    return _read_text_auto(p)


def _tables_docx(path: Path) -> list[list[list[str]]]:
    from docx import Document

    doc = Document(str(path))
    result: list[list[list[str]]] = []
    for table in doc.tables:
        result.append([[cell.text.strip() for cell in row.cells] for row in table.rows])
    return result


def _tables_pdf(path: Path) -> list[list[list[str]]]:
    import pdfplumber

    result: list[list[list[str]]] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            for table in page.extract_tables() or []:
                result.append([[cell or "" for cell in row] for row in table])
    return result


def extract_tables(path: str | Path) -> list[list[list[str]]]:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"Файл не найден: {p}")
    fmt = detect_format(p)
    if fmt == "docx":
        return _tables_docx(p)
    if fmt == "pdf":
        return _tables_pdf(p)
    return []


def read_xlsx_sheet(
    path: str | Path,
    sheet_name: str | None = None,
) -> list[list[str]]:
    from openpyxl import load_workbook

    p = Path(path)
    wb = load_workbook(str(p), data_only=True, read_only=True)
    target = sheet_name or wb.sheetnames[0]
    ws = wb[target]
    rows = [
        ["" if v is None else str(v) for v in row]
        for row in ws.iter_rows(values_only=True)
    ]
    wb.close()
    return rows


def _summarize_docx(path: Path) -> dict:
    from docx import Document

    doc = Document(str(path))
    cp = doc.core_properties
    return {
        "format": "docx",
        "size_bytes": path.stat().st_size,
        "author": cp.author or None,
        "created": str(cp.created) if cp.created else None,
        "modified": str(cp.modified) if cp.modified else None,
        "paragraphs": len(doc.paragraphs),
        "tables": len(doc.tables),
    }


def _summarize_pdf(path: Path) -> dict:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    meta = reader.metadata or {}
    return {
        "format": "pdf",
        "size_bytes": path.stat().st_size,
        "pages": len(reader.pages),
        "author": meta.get("/Author"),
        "created": str(meta.get("/CreationDate")) if meta.get("/CreationDate") else None,
        "modified": str(meta.get("/ModDate")) if meta.get("/ModDate") else None,
    }


def _summarize_pptx(path: Path) -> dict:
    from pptx import Presentation

    pres = Presentation(str(path))
    cp = pres.core_properties
    return {
        "format": "pptx",
        "size_bytes": path.stat().st_size,
        "slides": len(pres.slides),
        "author": cp.author or None,
        "created": str(cp.created) if cp.created else None,
        "modified": str(cp.modified) if cp.modified else None,
    }


def _summarize_xlsx(path: Path) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True)
    out = {
        "format": "xlsx",
        "size_bytes": path.stat().st_size,
        "sheets": list(wb.sheetnames),
    }
    wb.close()
    return out


def _summarize_xls(path: Path) -> dict:
    import xlrd

    book = xlrd.open_workbook(str(path))
    return {
        "format": "xls",
        "size_bytes": path.stat().st_size,
        "sheets": book.sheet_names(),
    }


def summarize(path: str | Path, *, preview_chars: int = 500) -> dict:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"Файл не найден: {p}")
    fmt = detect_format(p)
    if fmt == "docx":
        try:
            info = _summarize_docx(p)
        except Exception as e:
            info = {"format": "docx", "size_bytes": p.stat().st_size, "metadata_error": str(e)}
    elif fmt == "pdf":
        try:
            info = _summarize_pdf(p)
        except Exception as e:
            info = {"format": "pdf", "size_bytes": p.stat().st_size, "metadata_error": str(e)}
    elif fmt == "pptx":
        try:
            info = _summarize_pptx(p)
        except Exception as e:
            info = {"format": "pptx", "size_bytes": p.stat().st_size, "metadata_error": str(e)}
    elif fmt == "xlsx":
        try:
            info = _summarize_xlsx(p)
        except Exception as e:
            info = {"format": "xlsx", "size_bytes": p.stat().st_size, "metadata_error": str(e)}
    elif fmt == "xls":
        try:
            info = _summarize_xls(p)
        except Exception as e:
            info = {"format": "xls", "size_bytes": p.stat().st_size, "metadata_error": str(e)}
    else:
        info = {"format": fmt or "unknown", "size_bytes": p.stat().st_size}
    try:
        text = extract_text(p)
    except Exception as e:
        text = ""
        info["extract_error"] = str(e)
    info["preview"] = text[:preview_chars]
    return info
