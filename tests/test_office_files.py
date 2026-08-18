from __future__ import annotations

from datetime import datetime
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from workspace.utils import office_files as of


@pytest.fixture
def docx_path(tmp_path: Path) -> Path:
    p = tmp_path / "report.docx"
    doc = Document()
    doc.add_heading("Отчёт №42", level=1)
    doc.add_paragraph("Проверка проведена в 2026 году.")
    doc.add_paragraph("Выявлено 12 нарушений.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    doc.core_properties.author = "Иванов И.И."
    doc.core_properties.created = datetime(2026, 1, 15)
    doc.save(str(p))
    return p


@pytest.fixture
def xlsx_path(tmp_path: Path) -> Path:
    p = tmp_path / "data.xlsx"
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Нарушения"
    ws1["A1"] = "id"
    ws1["B1"] = "type"
    ws1["A2"] = 1
    ws1["B2"] = "финансовые"
    wb.create_sheet("Сводка")
    wb.save(str(p))
    return p


@pytest.fixture
def pptx_path(tmp_path: Path) -> Path:
    p = tmp_path / "deck.pptx"
    prs = Presentation()
    for title in ("Введение", "Результаты"):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
    prs.core_properties.author = "Петров"
    prs.save(str(p))
    return p


@pytest.fixture
def txt_path(tmp_path: Path) -> Path:
    p = tmp_path / "note.txt"
    p.write_text("Привет, мир!\nВторая строка.", encoding="utf-8")
    return p


@pytest.fixture
def csv_path(tmp_path: Path) -> Path:
    p = tmp_path / "data.csv"
    p.write_text("id,name\n1,Alpha\n2,Beta\n3,Gamma", encoding="utf-8")
    return p


def test_detect_format_by_extension():
    assert of.detect_format("a.docx") == "docx"
    assert of.detect_format("a.PDF") == "pdf"
    assert of.detect_format("a.Xlsx") == "xlsx"
    assert of.detect_format("noext") == ""


def test_extract_text_docx(docx_path: Path):
    text = of.extract_text(docx_path)
    assert "Отчёт №42" in text
    assert "12 нарушений" in text


def test_extract_tables_docx(docx_path: Path):
    tables = of.extract_tables(docx_path)
    assert len(tables) == 1
    assert tables[0][0] == ["A", "B"]
    assert tables[0][1] == ["1", "2"]


def test_summarize_docx(docx_path: Path):
    info = of.summarize(docx_path)
    assert info["format"] == "docx"
    assert info["author"] == "Иванов И.И."
    assert info["tables"] == 1
    assert info["paragraphs"] >= 3
    assert "preview" in info
    assert "Отчёт" in info["preview"]


def test_extract_text_xlsx(xlsx_path: Path):
    text = of.extract_text(xlsx_path)
    assert "Нарушения" in text
    assert "финансовые" in text


def test_summarize_xlsx(xlsx_path: Path):
    info = of.summarize(xlsx_path)
    assert info["format"] == "xlsx"
    assert "Нарушения" in info["sheets"]
    assert "Сводка" in info["sheets"]


def test_read_xlsx_sheet(xlsx_path: Path):
    rows = of.read_xlsx_sheet(xlsx_path, "Нарушения")
    assert rows[0] == ["id", "type"]
    assert rows[1] == ["1", "финансовые"]


def test_extract_text_pptx(pptx_path: Path):
    text = of.extract_text(pptx_path)
    assert "Введение" in text
    assert "Результаты" in text
    assert text.count("--- слайд") == 2


def test_summarize_pptx(pptx_path: Path):
    info = of.summarize(pptx_path)
    assert info["format"] == "pptx"
    assert info["slides"] == 2
    assert info["author"] == "Петров"


def test_extract_text_txt_utf8(txt_path: Path):
    text = of.extract_text(txt_path)
    assert "Привет" in text
    assert "Вторая строка" in text


def test_extract_text_txt_cp1251(tmp_path: Path):
    p = tmp_path / "cp1251.txt"
    p.write_bytes("Примечание: архив.".encode("cp1251"))
    text = of.extract_text(p)
    assert "Примечание" in text
    assert "архив" in text


def test_extract_text_csv(csv_path: Path):
    text = of.extract_text(csv_path)
    assert "Alpha" in text
    assert "Gamma" in text
    assert text.startswith("id | name")


def test_extract_text_missing_raises(tmp_path: Path):
    with pytest.raises(FileNotFoundError):
        of.extract_text(tmp_path / "missing.docx")


def test_summarize_missing_raises(tmp_path: Path):
    with pytest.raises(FileNotFoundError):
        of.summarize(tmp_path / "missing.pdf")


def test_extract_text_no_extension(tmp_path: Path):
    p = tmp_path / "noext"
    p.write_text("plain text", encoding="utf-8")
    assert of.extract_text(p) == "plain text"


def test_extract_tables_unknown_format(csv_path: Path):
    assert of.extract_tables(csv_path) == []


def test_extract_text_broken_docx_raises(tmp_path: Path):
    p = tmp_path / "bad.docx"
    p.write_bytes(b"not a docx")
    with pytest.raises(Exception):
        of.extract_text(p)


def test_extract_text_broken_xlsx_raises(tmp_path: Path):
    p = tmp_path / "bad.xlsx"
    p.write_bytes(b"garbage")
    with pytest.raises(Exception):
        of.extract_text(p)


def test_summarize_broken_docx_returns_metadata_error(tmp_path: Path):
    p = tmp_path / "bad.docx"
    p.write_bytes(b"corrupted")
    info = of.summarize(p)
    assert info["format"] == "docx"
    assert info.get("metadata_error") or info.get("extract_error")
